import os
import uuid
import zipfile
from io import BytesIO

from django.conf import settings
from django.http import JsonResponse, FileResponse, Http404
from rest_framework.decorators import api_view
import json

from pypdf import PdfReader, PdfWriter
import pdfplumber
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from fpdf import FPDF
import subprocess
try:
    import pytesseract
except ImportError:
    pytesseract = None
from reportlab.pdfgen import canvas
from reportlab.lib import pagesizes
from reportlab.lib import colors


def get_job_dir(job_id: str) -> str:
    job_dir = os.path.join(settings.MEDIA_ROOT, "jobs", job_id)
    os.makedirs(job_dir, exist_ok=True)
    return job_dir


def parse_ranges(range_text: str, total_pages: int):
    """
    Convert "1-3,5,8-10" into list of (start_idx, end_idx) 0-based inclusive
    """
    out = []
    s = (range_text or "").replace(" ", "")
    if not s:
        return out

    parts = s.split(",")
    for part in parts:
        if not part:
            continue
        if "-" in part:
            a, b = part.split("-", 1)
            start = int(a)
            end = int(b)
        else:
            start = int(part)
            end = int(part)

        # clamp to [1, total_pages]
        start = max(1, min(start, total_pages))
        end = max(1, min(end, total_pages))

        if end < start:
            start, end = end, start

        out.append((start - 1, end - 1))
    return out


def write_pdf_subset(reader: PdfReader, start_idx: int, end_idx: int, out_path: str):
    writer = PdfWriter()
    for i in range(start_idx, end_idx + 1):
        writer.add_page(reader.pages[i])
    with open(out_path, "wb") as f:
        writer.write(f)


@api_view(["POST"])
def split_prepare(request):
    """
    POST multipart:
      file: PDF
      mode: range|every|single
      range: (if range)
      every: (if every)
      base_name: (optional)
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)

    pdf_file = request.FILES["file"]
    mode = request.POST.get("mode", "range")
    range_text = request.POST.get("range", "")
    every = request.POST.get("every", "2")
    base_name = request.POST.get("base_name", "split").strip() or "split"

    try:
        reader = PdfReader(pdf_file)
        total = len(reader.pages)
        if total == 0:
            return JsonResponse({"error": "PDF has 0 pages"}, status=400)

        job_id = str(uuid.uuid4())
        job_dir = get_job_dir(job_id)

        results = []

        if mode == "single":
            # one PDF per page
            for i in range(total):
                out_name = f"{base_name}_page_{i+1}.pdf"
                out_path = os.path.join(job_dir, out_name)
                write_pdf_subset(reader, i, i, out_path)
                results.append({
                    "name": out_name,
                    "size": os.path.getsize(out_path),
                    "url": f"/api/jobs/{job_id}/file/{out_name}/",
                })

        elif mode == "every":
            n = int(every)
            if n <= 0:
                return JsonResponse({"error": "Every must be >= 1"}, status=400)

            chunk = 1
            for start in range(0, total, n):
                end = min(total - 1, start + n - 1)
                out_name = f"{base_name}_part_{chunk}.pdf"
                out_path = os.path.join(job_dir, out_name)
                write_pdf_subset(reader, start, end, out_path)
                results.append({
                    "name": out_name,
                    "size": os.path.getsize(out_path),
                    "url": f"/api/jobs/{job_id}/file/{out_name}/",
                })
                chunk += 1

        else:
            # range mode
            ranges = parse_ranges(range_text, total)
            if not ranges:
                return JsonResponse({"error": "Invalid or empty range"}, status=400)

            part = 1
            for (s, e) in ranges:
                out_name = f"{base_name}_range_{part}.pdf"
                out_path = os.path.join(job_dir, out_name)
                write_pdf_subset(reader, s, e, out_path)
                results.append({
                    "name": out_name,
                    "size": os.path.getsize(out_path),
                    "url": f"/api/jobs/{job_id}/file/{out_name}/",
                })
                part += 1

        return JsonResponse({
            "jobId": job_id,
            "results": results,
            "zipUrl": f"/api/jobs/{job_id}/zip/",
        })

    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


@api_view(["GET"])
def download_job_file(request, job_id, filename):
    job_dir = get_job_dir(job_id)
    safe = os.path.basename(filename)
    abs_path = os.path.join(job_dir, safe)
    if not os.path.exists(abs_path):
        raise Http404("File not found")

    return FileResponse(open(abs_path, "rb"), as_attachment=True, filename=safe)


@api_view(["GET"])
def download_job_zip(request, job_id):
    job_dir = get_job_dir(job_id)
    if not os.path.isdir(job_dir):
        raise Http404("Job not found")

    pdfs = sorted([f for f in os.listdir(job_dir) if f.lower().endswith(".pdf")])
    if not pdfs:
        return JsonResponse({"error": "No PDFs in this job"}, status=404)

    mem = BytesIO()
    with zipfile.ZipFile(mem, "w", zipfile.ZIP_DEFLATED) as zf:
        for fn in pdfs:
            zf.write(os.path.join(job_dir, fn), arcname=fn)

    mem.seek(0)
    return FileResponse(mem, as_attachment=True, filename=f"{job_id}_split.zip")


# --- New Features ---

from .services import compress_pdf as service_compress_pdf

@api_view(["POST"])
def compress_pdf_view(request):
    """
    POST multipart:
      file: PDF
      level: recommended|extreme|strong
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)
    
    pdf_file = request.FILES["file"]
    level = request.POST.get("level", "recommended")
    
    try:
        # Use service logic
        out_buffer, filename = service_compress_pdf(pdf_file, level)
        return FileResponse(out_buffer, as_attachment=True, filename=filename)
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


@api_view(["POST"])
def merge_pdfs(request):
    """
    POST multipart:
      files: list of PDF files
    """
    if "files" not in request.FILES:
        return JsonResponse({"error": "No files uploaded"}, status=400)
    
    files = request.FILES.getlist("files")
    if len(files) < 2:
        return JsonResponse({"error": "At least 2 files required for merge"}, status=400)

    try:
        merger = PdfWriter()
        for f in files:
            reader = PdfReader(f)
            for page in reader.pages:
                merger.add_page(page)
        
        out = BytesIO()
        merger.write(out)
        out.seek(0)
        return FileResponse(out, as_attachment=True, filename="merged.pdf")
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


@api_view(["POST"])
def rotate_pdf(request):
    """
    POST multipart:
      file: PDF
      angle: 90|180|270
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)
    
    pdf_file = request.FILES["file"]
    angle = int(request.POST.get("angle", 90))
    
    try:
        reader = PdfReader(pdf_file)
        writer = PdfWriter()
        
        for page in reader.pages:
            page.rotate(angle)
            writer.add_page(page)
            
        out = BytesIO()
        writer.write(out)
        out.seek(0)
        return FileResponse(out, as_attachment=True, filename="rotated.pdf")
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


@api_view(["POST"])
def delete_pages(request):
    """
    POST multipart:
      file: PDF
      pages: "1,3,5" (1-based)
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)
    
    pdf_file = request.FILES["file"]
    pages_to_delete_str = request.POST.get("pages", "")
    
    try:
        reader = PdfReader(pdf_file)
        total = len(reader.pages)
        
        # Parse logic similar to split but for EXCLUSION
        to_delete = set()
        for part in pages_to_delete_str.split(","):
            part = part.strip()
            if not part: continue
            if "-" in part:
                a, b = part.split("-")
                for x in range(int(a), int(b) + 1):
                    to_delete.add(x)
            else:
                to_delete.add(int(part))
        
        writer = PdfWriter()
        for i in range(total):
            page_num = i + 1
            if page_num not in to_delete:
                writer.add_page(reader.pages[i])
        
        out = BytesIO()
        writer.write(out)
        out.seek(0)
        return FileResponse(out, as_attachment=True, filename="deleted_pages.pdf")
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


@api_view(["POST"])
def extract_pages(request):
    """
    POST multipart:
      file: PDF
      pages: "1,3-5" (1-based)
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)
    
    pdf_file = request.FILES["file"]
    pages_str = request.POST.get("pages", "")
    
    try:
        reader = PdfReader(pdf_file)
        total = len(reader.pages)
        ranges = parse_ranges(pages_str, total) # Reuse existing logic
        
        writer = PdfWriter()
        for (start, end) in ranges:
            # ranges are 0-based inclusive
            for i in range(start, end + 1):
                writer.add_page(reader.pages[i])
                
        out = BytesIO()
        writer.write(out)
        out.seek(0)
        return FileResponse(out, as_attachment=True, filename="extracted.pdf")
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


@api_view(["POST"])
def protect_pdf(request):
    """
    POST multipart:
      file: PDF
      password: str
      permissions: comma-separated (print, copy, modify, annotate, extract)
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)

    pdf_file = request.FILES["file"]
    password = request.POST.get("password", "")
    permissions_arg = request.POST.get("permissions", "")

    try:
        import pikepdf

        pdf_bytes = pdf_file.read()
        pdf = pikepdf.open(BytesIO(pdf_bytes))

        # Build permissions object
        perms = [p.strip() for p in permissions_arg.split(",") if p.strip()]
        allow = pikepdf.Permissions(
            print_lowres=("print" in perms or not permissions_arg),
            print_highres=("print" in perms or not permissions_arg),
            modify_other=("modify" in perms or not permissions_arg),
            modify_annotation=("annotate" in perms or not permissions_arg),
            extract=("copy" in perms or "extract" in perms or not permissions_arg),
            assemble=("modify" in perms or not permissions_arg),
            form_filling=("annotate" in perms or not permissions_arg),
            accessibility=True,
        )

        out = BytesIO()
        # Use same password for user and owner if only one provided
        owner_pass = password if password else "owner_default_123"
        pdf.save(
            out,
            encryption=pikepdf.Encryption(
                user=password,
                owner=owner_pass,
                allow=allow,
            )
        )
        out.seek(0)
        return FileResponse(out, as_attachment=True, filename="protected.pdf")
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)



@api_view(["POST"])
def unlock_pdf(request):
    """
    POST multipart:
      file: PDF
      password: str
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)
    
    pdf_file = request.FILES["file"]
    password = request.POST.get("password", "")
    
    try:
        reader = PdfReader(pdf_file)
        if reader.is_encrypted:
            reader.decrypt(password)
        
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
            
        out = BytesIO()
        writer.write(out)
        out.seek(0)
        return FileResponse(out, as_attachment=True, filename="unlocked.pdf")
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


@api_view(["POST"])
def flatten_pdf(request):
    """
    POST multipart:
      file: PDF
    Flatten annotations and form fields.
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)
    
    pdf_file = request.FILES["file"]
    
    try:
        import fitz  # PyMuPDF

        pdf_bytes = pdf_file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")

        # Flatten form fields (widgets) at document level
        try:
            doc.flatten_form()
        except AttributeError:
            pass  # Older PyMuPDF versions may not have this

        # Flatten annotations (comments, highlights, etc.) at page level
        for page in doc:
            try:
                page.clean_contents()
            except Exception:
                pass
            try:
                page.flatten_annots()
            except Exception:
                pass

        out = BytesIO()
        doc.save(out)
        out.seek(0)
        return FileResponse(out, as_attachment=True, filename="flattened.pdf")
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)

from pdf2docx import Converter
import img2pdf

@api_view(["POST"])
def pdf_to_word(request):
    """
    POST multipart:
      file: PDF
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)
    
    pdf_file = request.FILES["file"]
    job_id = str(uuid.uuid4())
    job_dir = get_job_dir(job_id)
    
    input_path = os.path.join(job_dir, "input.pdf")
    output_path = os.path.join(job_dir, "converted.docx")
    
    try:
        # Save temp file
        with open(input_path, "wb") as f:
            for chunk in pdf_file.chunks():
                f.write(chunk)
        
        cv = Converter(input_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()
        
        return FileResponse(open(output_path, "rb"), as_attachment=True, filename="converted.docx")
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


@api_view(["POST"])
def pdf_to_image(request):
    """
    POST multipart:
      file: PDF
    Returns ZIP of images
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)
    
    pdf_file = request.FILES["file"]
    job_id = str(uuid.uuid4())
    job_dir = get_job_dir(job_id)
    
    try:
        import fitz  # PyMuPDF
        
        # Read from memory
        pdf_bytes = pdf_file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        
        image_files = []
        for i in range(len(doc)):
            page = doc.load_page(i)
            # Render high quality
            pix = page.get_pixmap(dpi=150)
            img_name = f"page_{i+1}.png"
            img_path = os.path.join(job_dir, img_name)
            pix.save(img_path)
            image_files.append(img_name)
            
        # Zip them
        mem = BytesIO()
        with zipfile.ZipFile(mem, "w", zipfile.ZIP_DEFLATED) as zf:
            for fn in image_files:
                zf.write(os.path.join(job_dir, fn), arcname=fn)
        
        mem.seek(0)
        return FileResponse(mem, as_attachment=True, filename="pdf_images.zip")
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


@api_view(["POST"])
def image_to_pdf(request):
    """
    POST multipart:
      files: list of images
    """
    if "files" not in request.FILES:
        return JsonResponse({"error": "No files uploaded"}, status=400)
    
    files = request.FILES.getlist("files")
    if not files:
        return JsonResponse({"error": "No image files found"}, status=400)
    
    try:
        # img2pdf requires files or bytes
        # We'll read all into a list of bytes
        images_bytes = []
        for f in files:
            f.seek(0)
            images_bytes.append(f.read())
            
        pdf_bytes = img2pdf.convert(images_bytes)
        
        return FileResponse(BytesIO(pdf_bytes), as_attachment=True, filename="images.pdf")
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)

@api_view(["POST"])
def auto_rename_pdf(request):
    """
    POST multipart:
      file: PDF
    Returns the same file with a new filename based on metadata or content.
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)
    
    pdf_file = request.FILES["file"]
    
    try:
        reader = PdfReader(pdf_file)
        new_name = "renamed.pdf"
        
        # 1. Try Metadata Title
        if reader.metadata and reader.metadata.title:
            title = reader.metadata.title.strip()
            if title:
                # Basic sanitation
                safe_title = "".join(c for c in title if c.isalnum() or c in (" ", "-", "_")).strip()
                if safe_title:
                    new_name = f"{safe_title}.pdf"
        
        # 2. If no title, try text from first page
        if new_name == "renamed.pdf" and len(reader.pages) > 0:
            text = reader.pages[0].extract_text()
            if text:
                lines = [l.strip() for l in text.splitlines() if l.strip()]
                if lines:
                    candidate = lines[0]
                    # Take first 50 chars, sanitize
                    safe_candidate = "".join(c for c in candidate[:50] if c.isalnum() or c in (" ", "-", "_")).strip()
                    if safe_candidate:
                        new_name = f"{safe_candidate}.pdf"

        # Return the original file content with new name
        # We need to rewind the file ptr
        pdf_file.seek(0)
        return FileResponse(pdf_file, as_attachment=True, filename=new_name)

    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


@api_view(["POST"])
def reorder_pages(request):
    """
    POST multipart:
      file: PDF
      order: "1,3,2" (1-based)
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)
    
    pdf_file = request.FILES["file"]
    order_str = request.POST.get("order", "")
    
    try:
        reader = PdfReader(pdf_file)
        total = len(reader.pages)
        writer = PdfWriter()
        
        # Parse order similar to ranges but we want explicit order
        parts = order_str.split(",")
        for part in parts:
            part = part.strip()
            if not part: continue
            
            if "-" in part:
                 # Range handling if user inputs 1-3
                 a, b = part.split("-")
                 start, end = int(a), int(b)
                 # Clamp
                 start = max(1, min(start, total))
                 end = max(1, min(end, total))
                 if start > end: start, end = end, start
                 for i in range(start, end + 1):
                     writer.add_page(reader.pages[i-1])
            else:
                 idx = int(part)
                 if 1 <= idx <= total:
                     writer.add_page(reader.pages[idx-1])

        out = BytesIO()
        writer.write(out)
        out.seek(0)
        return FileResponse(out, as_attachment=True, filename="reordered.pdf")
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)

# --- View & Edit Features ---

@api_view(["POST"])
def number_pages(request):
    """
    POST multipart:
      file: PDF
      position: bottom-center|bottom-right|bottom-left|top-center|...
      style: "1", "Page 1", "Page 1 of N"
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)
    
    pdf_file = request.FILES["file"]
    position = request.POST.get("position", "bottom-center")
    style = request.POST.get("style", "1")
    
    try:
        reader = PdfReader(pdf_file)
        writer = PdfWriter()
        total_pages = len(reader.pages)
        
        for i, page in enumerate(reader.pages):
            # Create a PDF with just the number
            packet = BytesIO()
            
            # Get page size in points (1/72 inch)
            width = float(page.mediabox.width)
            height = float(page.mediabox.height)
            
            c = canvas.Canvas(packet, pagesize=(width, height))
            
            # Determine text
            page_num = i + 1
            if style == "Page 1 of N":
                text = f"Page {page_num} of {total_pages}"
            elif style == "Page 1":
                text = f"Page {page_num}"
            else:
                text = f"{page_num}"
            
            c.setFont("Helvetica", 12)
            
            # Position logic (margin 20pts approx 7mm)
            margin = 20
            text_width = c.stringWidth(text, "Helvetica", 12)
            
            if "bottom" in position:
                y = margin
            elif "top" in position:
                y = height - margin - 10 # adjust for font height
            else:
                y = margin
                
            if "left" in position:
                x = margin
            elif "right" in position:
                x = width - margin - text_width
            else: # center
                x = (width - text_width) / 2
            
            c.drawString(x, y, text)
            c.save()
            
            packet.seek(0)
            number_pdf = PdfReader(packet)
            
            # Merge
            page.merge_page(number_pdf.pages[0])
            writer.add_page(page)
            
        out = BytesIO()
        writer.write(out)
        out.seek(0)
        return FileResponse(out, as_attachment=True, filename="numbered.pdf")
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


@api_view(["POST"])
def crop_pdf(request):
    """
    POST multipart:
      file: PDF
      box: json string {x, y, width, height} or "margin"
      # For simplicity, let's assume we crop all pages to a margin or specific box relative to page size?
      # Actually, let's implement a simple margin crop or fixed box if provided.
      # If "box" is provided as "10,10,200,200" (x,y,w,h)
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)
    
    pdf_file = request.FILES["file"]
    box_str = request.POST.get("box", "")
    
    try:
        reader = PdfReader(pdf_file)
        writer = PdfWriter()
        
        for page in reader.pages:
            # If box provided
            if box_str:
                try:
                    # x,y,w,h
                    parts = [float(p) for p in box_str.split(",")]
                    if len(parts) == 4:
                        x, y, w, h = parts
                        # pypdf RectangleObject: [left, bottom, right, top]
                        # Assume x,y are from bottom-left or top-left? 
                        # PDF coords are usually bottom-left. 
                        # Let's assume input is x,y (bottom-left), w, h
                        page.mediabox.lower_left = (x, y)
                        page.mediabox.upper_right = (x + w, y + h)
                except:
                    pass # ignore invalid box, just copy page
            
            writer.add_page(page)
            
        out = BytesIO()
        writer.write(out)
        out.seek(0)
        return FileResponse(out, as_attachment=True, filename="cropped.pdf")
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


@api_view(["POST"])
def watermark_pdf(request):
    """
    POST multipart:
      file: PDF
      text: str (optional)
      image: file (optional)
      opacity: float (0.1 to 1.0)
      position: center|...
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)
    
    pdf_file = request.FILES["file"]
    watermark_text = request.POST.get("text", "")
    watermark_image = request.FILES.get("image")
    
    try:
        reader = PdfReader(pdf_file)
        writer = PdfWriter()
        
        # Create watermark PDF
        
        for page in reader.pages:
            width = float(page.mediabox.width)
            height = float(page.mediabox.height)
            
            packet = BytesIO()
            c = canvas.Canvas(packet, pagesize=(width, height))
            
            # Draw watermark
            if watermark_text:
                c.setFont("Helvetica-Bold", 48)
                c.setFillColorRGB(0.5, 0.5, 0.5, alpha=0.3)
                c.saveState()
                c.translate(width/2, height/2)
                c.rotate(45)
                c.drawCentredString(0, 0, watermark_text)
                c.restoreState()
            elif watermark_image:
                 # TODO: Implement image watermark using c.drawImage
                 pass
                 
            c.save()
            packet.seek(0)
            overlay = PdfReader(packet)
            if len(overlay.pages) > 0:
                page.merge_page(overlay.pages[0])
            
            writer.add_page(page)
            
        out = BytesIO()
        writer.write(out)
        out.seek(0)
        return FileResponse(out, as_attachment=True, filename="watermarked.pdf")
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


@api_view(["POST"])
def form_fill(request):
    """
    POST multipart:
      file: PDF
      fields: json string { "field_name": "value" }
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)
        
    pdf_file = request.FILES["file"]
    fields_str = request.POST.get("fields", "{}")
    
    try:
        fields = json.loads(fields_str)
        reader = PdfReader(pdf_file)
        writer = PdfWriter()
        
        # Copy pages and update fields
        writer.append_pages_from_reader(reader)
        writer.update_page_form_field_values(
            writer.pages[0], fields
        )
        
        # For multiple pages, we might need to iterate.
        # simple update_page_form_field_values might only work if we used clone_reader_document_root or similar.
        # Actually pypdf's form filling is bit complex. 
        # But let's try the standard way.
        
        out = BytesIO()
        writer.write(out)
        out.seek(0)
        return FileResponse(out, as_attachment=True, filename="filled.pdf")
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


@api_view(["POST"])
def pdf_to_excel(request):
    """
    POST multipart:
      file: PDF
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)
    
    pdf_file = request.FILES["file"]
    job_id = str(uuid.uuid4())
    job_dir = get_job_dir(job_id)
    
    input_path = os.path.join(job_dir, "input.pdf")
    output_path = os.path.join(job_dir, "converted.xlsx")
    
    try:
        with open(input_path, "wb") as f:
            for chunk in pdf_file.chunks():
                f.write(chunk)
                
        # Basic extraction: extract tables from each page
        with pdfplumber.open(input_path) as pdf:
            all_tables = []
            for page in pdf.pages:
                tables = page.extract_tables()
                for table in tables:
                    df = pd.DataFrame(table)
                    all_tables.append(df)
            
            if not all_tables:
                 pd.DataFrame(["No tables found"]).to_excel(output_path, index=False, header=False)
            else:
                final_df = pd.concat(all_tables, ignore_index=True)
                final_df.to_excel(output_path, index=False, header=False)

        return FileResponse(open(output_path, "rb"), as_attachment=True, filename="converted.xlsx")
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


@api_view(["POST"])
def pdf_to_ppt(request):
    """
    POST multipart:
      file: PDF
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)

    pdf_file = request.FILES["file"]
    job_id = str(uuid.uuid4())
    job_dir = get_job_dir(job_id)

    input_path = os.path.join(job_dir, "input.pdf")
    output_path = os.path.join(job_dir, "converted.pptx")

    try:
        with open(input_path, "wb") as f:
            for chunk in pdf_file.chunks():
                f.write(chunk)
        
        import fitz # PyMuPDF
        
        doc = fitz.open(input_path)
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        
        for i in range(len(doc)):
            page = doc.load_page(i)
            pix = page.get_pixmap(dpi=150)
            img_name = f"page_{i}.png"
            img_path = os.path.join(job_dir, img_name)
            pix.save(img_path)
            
            blank_slide_layout = prs.slide_layouts[6] 
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            
        prs.save(output_path)
        
        return FileResponse(open(output_path, "rb"), as_attachment=True, filename="converted.pptx")
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


@api_view(["POST"])
def pdf_to_text(request):
    """
    POST multipart:
      file: PDF
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)

    pdf_file = request.FILES["file"]
    
    try:
        reader = PdfReader(pdf_file)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n\n"
            
        out = BytesIO()
        out.write(text.encode("utf-8"))
        out.seek(0)
        
        return FileResponse(out, as_attachment=True, filename="converted.txt")
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


@api_view(["POST"])
def word_to_pdf(request):
    """
    POST multipart:
      file: DOCX/DOC
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)

    docx_file = request.FILES["file"]
    job_id = str(uuid.uuid4())
    job_dir = get_job_dir(job_id)

    original_name = docx_file.name
    input_path = os.path.join(job_dir, original_name)
    output_path = os.path.join(job_dir, os.path.splitext(original_name)[0] + ".pdf")

    try:
        with open(input_path, "wb") as f:
            for chunk in docx_file.chunks():
                f.write(chunk)
        
        soffice = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        if os.path.exists(soffice):
             subprocess.run([soffice, "--headless", "--convert-to", "pdf", input_path, "--outdir", job_dir], check=True)
        else:
             # Fallback attempt using textutil for basic docx to html then... wait, textutil can convert docx to html.
             # but we need PDF. macOS 'cupsfilter' can convert text/html to pdf.
             # Let's try: textutil -convert html input.docx -stdout | cupsfilter -i text/html - > output.pdf
             # This is a bit hacky but might work on macOS for basic files.
             
             try:
                 # 1. Convert to HTML
                 html_path = os.path.join(job_dir, "temp.html")
                 subprocess.run(["textutil", "-convert", "html", input_path, "-output", html_path], check=True)
                 
                 # 2. Convert HTML to PDF using cupsfilter
                 with open(output_path, "wb") as pdf_out:
                     # cupsfilter writes to stdout
                     subprocess.run(["cupsfilter", "-i", "text/html", html_path], stdout=pdf_out, check=True)
             except Exception as sub_e:
                 return JsonResponse({"error": "LibreOffice not found and fallback failed: " + str(sub_e)}, status=501)

        if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
             return JsonResponse({"error": "Conversion failed to produce PDF output"}, status=500)
             
        return FileResponse(open(output_path, "rb"), as_attachment=True, filename=os.path.basename(output_path))
    except Exception as e:
        return JsonResponse({"error": f"Error: {str(e)}"}, status=500)


@api_view(["POST"])
def excel_to_pdf(request):
    """
    POST multipart:
      file: XLS/XLSX
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)

    f = request.FILES["file"]
    job_id = str(uuid.uuid4())
    job_dir = get_job_dir(job_id)
    input_path = os.path.join(job_dir, f.name)

    try:
        with open(input_path, "wb") as dest:
            for chunk in f.chunks():
                dest.write(chunk)
        
        soffice = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        if os.path.exists(soffice):
             subprocess.run([soffice, "--headless", "--convert-to", "pdf", input_path, "--outdir", job_dir], check=True)
             expected_out = os.path.splitext(input_path)[0] + ".pdf"
             return FileResponse(open(expected_out, "rb"), as_attachment=True, filename=os.path.basename(expected_out))
        else:
            return JsonResponse({"error": "LibreOffice not found. Excel to PDF not supported."}, status=501)
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


@api_view(["POST"])
def ppt_to_pdf(request):
    """
    POST multipart:
      file: PPT/PPTX
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)

    f = request.FILES["file"]
    job_id = str(uuid.uuid4())
    job_dir = get_job_dir(job_id)
    input_path = os.path.join(job_dir, f.name)

    try:
        with open(input_path, "wb") as dest:
            for chunk in f.chunks():
                dest.write(chunk)

        soffice = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        if os.path.exists(soffice):
             subprocess.run([soffice, "--headless", "--convert-to", "pdf", input_path, "--outdir", job_dir], check=True)
             expected_out = os.path.splitext(input_path)[0] + ".pdf"
             return FileResponse(open(expected_out, "rb"), as_attachment=True, filename=os.path.basename(expected_out))
        else:
            return JsonResponse({"error": "LibreOffice not found. PPT to PDF not supported."}, status=501)
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


@api_view(["POST"])
def text_to_pdf(request):
    """
    POST multipart:
      file: TXT
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)

    f = request.FILES["file"]
    
    try:
        content = f.read().decode("utf-8", errors="ignore")
        
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        pdf.multi_cell(0, 10, content)
        
        s = pdf.output(dest='S').encode('latin-1')
        return FileResponse(BytesIO(s), as_attachment=True, filename="converted.pdf")
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


@api_view(["POST"])
def ocr_pdf(request):
    """
    POST multipart:
      file: PDF (scanned)
    """
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)

    if pytesseract is None:
        return JsonResponse({"error": "pytesseract not installed on server"}, status=501)

    pdf_file = request.FILES["file"]
    job_id = str(uuid.uuid4())
    job_dir = get_job_dir(job_id)
    input_path = os.path.join(job_dir, "input.pdf")

    try:
        with open(input_path, "wb") as f:
            for chunk in pdf_file.chunks():
                f.write(chunk)
                
        import fitz
        doc = fitz.open(input_path)
        merger = PdfWriter()
        
        for i in range(len(doc)):
            page = doc.load_page(i)
            # Higher DPI for OCR
            pix = page.get_pixmap(dpi=300)
            img_path = os.path.join(job_dir, f"page_{i}.png")
            pix.save(img_path)
            
            try:
                pdf_bytes = pytesseract.image_to_pdf_or_hocr(img_path, extension='pdf')
                merger.append(BytesIO(pdf_bytes))
            except Exception as tess_err:
                print(f"OCR failed for page {i}: {tess_err}")
                return JsonResponse({"error": "Tesseract OCR binary not found or failed"}, status=500)

        out = BytesIO()
        merger.write(out)
        out.seek(0)
        
        return FileResponse(out, as_attachment=True, filename="ocr_result.pdf")
        
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)

@api_view(["POST"])
def sign_pdf(request):
    """
    POST multipart:
      file: PDF
      signature: Image (PNG/JPG)
      page: int (1-based, default 1)
      x: int (default 100)
      y: int (default 100)
    """
    if "file" not in request.FILES or "signature" not in request.FILES:
        return JsonResponse({"error": "PDF and Signature files are required"}, status=400)
    
    pdf_file = request.FILES["file"]
    sig_file = request.FILES["signature"]
    page_num = int(request.POST.get("page", 1)) - 1
    x = int(request.POST.get("x", 100))
    y = int(request.POST.get("y", 100))
    
    try:
        import fitz  # PyMuPDF
        
        # Open PDF
        pdf_bytes = pdf_file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        
        if page_num < 0 or page_num >= len(doc):
             return JsonResponse({"error": "Invalid page number"}, status=400)
             
        page = doc[page_num]
        
        # Open Signature Image
        sig_bytes = sig_file.read()
        
        # Insert Image
        # rect is (x0, y0, x1, y1)
        # We'll assume a default width/height if not specified, or scale it?
        # Let's say we want it to be 100x50 by default effectively? 
        # Actually fitz.insert_image takes kwarg 'stream'
        
        # Let's try to get image size first to keep aspect ratio if possible, or just standard box
        # For MVP, let's fix a box of width 150, height 75 at (x, y)
        w, h = 150, 75
        rect = fitz.Rect(x, y, x + w, y + h)
        
        page.insert_image(rect, stream=sig_bytes)
        
        out_buffer = BytesIO()
        doc.save(out_buffer)
        out_buffer.seek(0)
        
        return FileResponse(out_buffer, as_attachment=True, filename="signed.pdf")
        
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)
